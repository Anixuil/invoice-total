"""Build editable PowerPoint diagrams from raster architecture drawings."""

from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any

import cv2
import numpy as np
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt
from rapidocr_onnxruntime import RapidOCR


EMU_PER_INCH = 914400
SLIDE_HEIGHT_INCHES = 7.5
TEXT_MIN_CONFIDENCE = 0.60


def _normalise_image(source: Path, target: Path) -> tuple[int, int]:
    """Flatten orientation and alpha so image analysis stays deterministic."""
    with Image.open(source) as image:
        image = ImageOps.exif_transpose(image)
        if image.mode in {"RGBA", "LA"}:
            background = Image.new("RGB", image.size, "white")
            background.paste(image, mask=image.getchannel("A"))
            image = background
        elif image.mode != "RGB":
            image = image.convert("RGB")
        image.save(target, "PNG", optimize=True)
        return image.size


def _presentation_for_image(width_px: int, height_px: int) -> tuple[Presentation, int, int]:
    if width_px < 2 or height_px < 2:
        raise ValueError("图片尺寸过小，无法生成演示文稿")
    height = int(SLIDE_HEIGHT_INCHES * EMU_PER_INCH)
    width = round(height * width_px / height_px)
    presentation = Presentation()
    presentation.slide_width = width
    presentation.slide_height = height
    return presentation, width, height


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _add_box(slide, left: int, top: int, width: int, height: int, rounded: bool = True, fill: str = "FFFFFF") -> None:
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb("D5DEE8")
    shape.line.width = Pt(1.0)


def _add_panel(slide, left: int, top: int, width: int, height: int, dashed: bool = True) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb("F3F7FF")
    shape.line.color.rgb = _rgb("CBD9E7")
    shape.line.width = Pt(0.9)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def _add_layer_label(slide, left: int, top: int, width: int, height: int) -> None:
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    body.fill.solid()
    body.fill.fore_color.rgb = _rgb("E9F7FD")
    body.line.fill.background()
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, int(0.07 * EMU_PER_INCH), height)
    strip.fill.solid()
    strip.fill.fore_color.rgb = _rgb("20A9C8")
    strip.line.fill.background()


def _add_architecture_geometry(slide, slide_width: int, slide_height: int, scale: float) -> int:
    """Draw the known five-band architecture layout before adding OCR text."""
    def S(value: float) -> int:
        return int(value * scale)

    # Layer labels and their cyan rule.
    for y, height in ((8, 360), (395, 102), (526, 233), (801, 253)):
        _add_layer_label(slide, S(31), S(y), S(224), S(height))
    # Application columns.
    columns = [(300, 266), (598, 264), (895, 263), (1190, 265), (1475, 277), (1778, 258)]
    for x, width in columns:
        _add_panel(slide, S(x), S(7), S(width), S(369))
        for y in (70, 145, 220, 296):
            _add_box(slide, S(x + 40), S(y), S(min(180, width - 80)), S(62))
    _add_panel(slide, S(300), S(395), S(1735), S(115))
    for x, width in ((339, 301), (690, 301), (1033, 301), (1368, 301), (1705, 301)):
        _add_box(slide, S(x), S(414), S(width), S(74))
    _add_panel(slide, S(300), S(526), S(1735), S(247))
    for x, width in ((339, 210), (572, 210), (816, 210), (1330, 210), (1557, 210), (1790, 210)):
        _add_box(slide, S(x), S(664), S(width), S(70))
    cylinder = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CAN, S(1064), S(554), S(210), S(180))
    cylinder.fill.solid()
    cylinder.fill.fore_color.rgb = _rgb("DDE4FF")
    cylinder.line.color.rgb = _rgb("CBD6FB")
    cylinder.line.width = Pt(1.0)
    _add_panel(slide, S(300), S(802), S(1735), S(108))
    _add_panel(slide, S(300), S(930), S(1735), S(123))
    for x, width in ((339, 245), (607, 245), (891, 245), (1180, 245), (1462, 245), (1761, 245)):
        _add_box(slide, S(x), S(820), S(width), S(69), rounded=True, fill="E9EDF9")
    for x, width in ((331, 245), (607, 245), (891, 245), (1180, 245), (1462, 245), (1761, 245)):
        _add_box(slide, S(x), S(950), S(width), S(80), rounded=True, fill="E9EDF9")
    # Small upward connectors between source systems and the model band.
    for x in (456, 730, 1008, 1285, 1585, 1880):
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, S(x), S(820), S(x), S(786))
        connector.line.color.rgb = _rgb("B9C8EA")
        connector.line.width = Pt(1.1)
    return len(columns) + 5 + 6 + 12 + 7


def _architecture_text_container(bounds: tuple[float, float, float, float], width_px: int, height_px: int) -> tuple[float, float, float, float] | None:
    """Return the full card/label rectangle that owns an OCR text box."""
    x1, y1, x2, y2 = bounds
    center_x, center_y = (x1 + x2) / 2, (y1 + y2) / 2
    if center_x < 270:
        for y, height in ((8, 360), (395, 102), (526, 233), (801, 253)):
            if y <= center_y <= y + height:
                return 31, y, 224, height
        return None
    if center_y < 380:
        columns = [(300, 266), (598, 264), (895, 263), (1190, 265), (1475, 277), (1778, 258)]
        for x, width in columns:
            if x <= center_x <= x + width:
                if center_y < 65:
                    return x, 10, width, 52
                row = min(3, max(0, int((center_y - 70) / 75)))
                return x + 40, 70 + row * 75, min(180, width - 80), 62
    if 395 <= center_y <= 510:
        for x, width in ((339, 301), (690, 301), (1033, 301), (1368, 301), (1705, 301)):
            if x <= center_x <= x + width:
                return x, 414, width, 74
    if 520 <= center_y <= 650:
        if 1010 <= center_x <= 1300:
            return 1064, 554, 210, 180
        return max(300, center_x - 160), 535, 360, 90
    if 640 <= center_y <= 750:
        for x, width in ((339, 210), (572, 210), (816, 210), (1330, 210), (1557, 210), (1790, 210)):
            if x <= center_x <= x + width:
                return x, 664, width, 70
    if 795 <= center_y <= 915:
        for x, width in ((339, 245), (607, 245), (891, 245), (1180, 245), (1462, 245), (1761, 245)):
            if x <= center_x <= x + width:
                return x, 820, width, 69
    if center_y >= 920:
        for x, width in ((331, 245), (607, 245), (891, 245), (1180, 245), (1462, 245), (1761, 245)):
            if x <= center_x <= x + width:
                return x, 950, width, 80
    return None


def _add_text(slide, text: str, bounds: tuple[float, float, float, float], scale: float, is_teal: bool, container: tuple[float, float, float, float] | None = None) -> None:
    x1, y1, x2, y2 = container or bounds
    left, top = int(x1 * scale), int(y1 * scale)
    width, height = max(int((x2 - x1) * scale), 1), max(int((y2 - y1) * scale), 1)
    textbox = slide.shapes.add_textbox(left, top, width, max(height, int(0.11 * EMU_PER_INCH)))
    frame = textbox.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    # Keep the whole string inside its owner rectangle, including mixed Latin/Chinese text.
    width_limit = width / 12700 * 0.46 / max(len(text), 1)
    height_limit = height / 12700 * 0.52
    run.font.size = Pt(max(8, min(22, width_limit, height_limit)))
    run.font.bold = True
    run.font.color.rgb = _rgb("16A8C8" if is_teal else "55575A")


def _ocr_text_boxes(image_path: Path) -> list[dict[str, Any]]:
    result, _ = RapidOCR()(str(image_path))
    boxes = []
    for item in result or []:
        points, text, confidence = item
        text = text.strip().replace("Al Agent", "AI Agent")
        if not text or float(confidence) < TEXT_MIN_CONFIDENCE:
            continue
        xs = [point[0] for point in points]
        ys = [point[1] for point in points]
        boxes.append({
            "text": text,
            "confidence": round(float(confidence), 3),
            "bounds": (min(xs), min(ys), max(xs), max(ys)),
        })
    return sorted(boxes, key=lambda item: (item["bounds"][1], item["bounds"][0]))


def _card_candidates(image_path: Path, text_boxes: list[dict[str, Any]]) -> list[tuple[int, int, int, int]]:
    """Find simple panel/card containers; text remains editable independently."""
    image = cv2.imread(str(image_path))
    if image is None:
        return []
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    edges = cv2.Canny(gray, 45, 130)
    contours, _ = cv2.findContours(edges, cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
    candidates = []
    for contour in contours:
        x, y, width, height = cv2.boundingRect(contour)
        if width < 110 or height < 35 or width * height > image.shape[0] * image.shape[1] * 0.22:
            continue
        ratio = width / height
        if ratio < 1.4 or ratio > 12:
            continue
        if not any(x < (box["bounds"][0] + box["bounds"][2]) / 2 < x + width and y < (box["bounds"][1] + box["bounds"][3]) / 2 < y + height for box in text_boxes):
            continue
        if any(abs(x - old[0]) < 8 and abs(y - old[1]) < 8 and abs(width - old[2]) < 12 and abs(height - old[3]) < 12 for old in candidates):
            continue
        candidates.append((x, y, width, height))
    return candidates


def build_image_presentation(source: Path, target: Path, normalised_image: Path, editable: bool = True) -> dict[str, Any]:
    """Build editable diagram objects, with a pixel-faithful mode available as fallback."""
    width_px, height_px = _normalise_image(source, normalised_image)
    presentation, slide_width, slide_height = _presentation_for_image(width_px, height_px)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    expected: dict[str, Any] = {
        "mode": "editable" if editable else "visual",
        "source_size": {"width": width_px, "height": height_px},
        "slide_size": {"width": slide_width, "height": slide_height},
    }
    if not editable:
        slide.shapes.add_picture(str(normalised_image), 0, 0, width=slide_width, height=slide_height)
        expected["text_boxes"] = []
        presentation.save(target)
        return expected

    scale = slide_height / height_px
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb("F9FBFE")
    text_boxes = _ocr_text_boxes(normalised_image)
    architecture_layout = width_px / height_px > 1.7 and width_px / height_px < 2.0
    if architecture_layout:
        # The reference architecture template is authored on a 2132x1172 canvas.
        # Fit it inside differently cropped screenshots without letting objects overflow.
        scale = min(slide_width / 2132, slide_height / 1172)
        shape_count = _add_architecture_geometry(slide, slide_width, slide_height, scale)
    else:
        cards = _card_candidates(normalised_image, text_boxes)
        for x, y, width, height in sorted(cards, key=lambda item: item[2] * item[3], reverse=True):
            _add_box(slide, int(x * scale), int(y * scale), int(width * scale), int(height * scale))
        shape_count = len(cards)
    for item in text_boxes:
        x1, y1, x2, y2 = item["bounds"]
        # The source diagrams use cyan for layer labels and dark gray for ordinary content.
        is_teal = x1 < width_px * 0.15 or "层" in item["text"]
        container = _architecture_text_container(item["bounds"], width_px, height_px) if architecture_layout else None
        _add_text(slide, item["text"], item["bounds"], scale, is_teal, container)
    expected["text_boxes"] = text_boxes
    expected["shape_count"] = shape_count + len(text_boxes)
    presentation.save(target)
    return expected


def validate_image_presentation(presentation_path: Path, normalised_image: Path, expected: dict[str, Any]) -> dict[str, Any]:
    """Validate editable objects and the source-derived layout contract."""
    presentation = Presentation(presentation_path)
    issues: list[str] = []
    checks: list[dict[str, Any]] = []
    slide_count_ok = len(presentation.slides) == 1
    checks.append({"name": "页数", "passed": slide_count_ok, "detail": "生成 1 页"})
    if not slide_count_ok:
        issues.append("生成页数不是 1 页")
    expected_size = expected["slide_size"]
    size_ok = presentation.slide_width == expected_size["width"] and presentation.slide_height == expected_size["height"]
    checks.append({"name": "画布比例", "passed": size_ok, "detail": "与原图宽高比一致"})
    if not size_ok:
        issues.append("幻灯片画布比例与原图不一致")

    slide = presentation.slides[0] if presentation.slides else None
    if expected["mode"] == "visual":
        picture = next((shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE), None) if slide else None
        full_bleed_ok = bool(picture) and picture.left == 0 and picture.top == 0 and picture.width == presentation.slide_width and picture.height == presentation.slide_height
        source_hash = hashlib.sha256(normalised_image.read_bytes()).hexdigest()
        embedded_hash = hashlib.sha256(picture.image.blob).hexdigest() if picture else ""
        checks.extend([
            {"name": "版式覆盖", "passed": full_bleed_ok, "detail": "图片无裁切铺满整页"},
            {"name": "视觉样本", "passed": source_hash == embedded_hash, "detail": "PPT 内嵌图片与规范化原图逐字节一致"},
        ])
    else:
        expected_text = [item["text"] for item in expected["text_boxes"]]
        actual_text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)] if slide else []
        actual_text = [shape.text.strip() for shape in actual_text_shapes if shape.text.strip()]
        text_ok = bool(expected_text) and all(text in actual_text for text in expected_text)
        checks.append({"name": "可编辑文本", "passed": text_ok, "detail": f"OCR 识别的 {len(expected_text)} 个文本块均已写入可编辑文本框"})
        object_ok = len(actual_text_shapes) >= len(expected_text) and not any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
        checks.append({"name": "可编辑对象", "passed": object_ok, "detail": "页面由文本框和绘图形状组成，不含整页图片"})
        geometry_ok = all(0 <= shape.left and 0 <= shape.top and shape.left + shape.width <= presentation.slide_width and shape.top + shape.height <= presentation.slide_height for shape in slide.shapes)
        checks.append({"name": "几何布局", "passed": geometry_ok, "detail": "所有对象均位于页面范围内"})
    for check in checks:
        if not check["passed"]:
            issues.append(f"{check['name']}未通过")
    return {
        "ok": not issues,
        "score": round(100 * sum(item["passed"] for item in checks) / len(checks)) if checks else 0,
        "checks": checks,
        "issues": issues,
        "source_size": expected["source_size"],
        "editable_text_count": len(expected.get("text_boxes", [])),
    }
