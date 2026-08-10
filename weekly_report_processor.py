#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""部门项目周报的 PPT 审核、模板化合并和周例会文档生成。"""

from __future__ import annotations

from copy import deepcopy
from datetime import date, datetime, timedelta
from difflib import SequenceMatcher
from functools import lru_cache
from io import BytesIO
from pathlib import Path
import re
from typing import Any, Callable

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Pt


ProgressCallback = Callable[[dict[str, Any]], None] | None
BASE = Path(__file__).resolve().parent
PPT_TEMPLATE = BASE / "templates" / "weekly_report_template.pptx"
DOCX_TEMPLATE = BASE / "templates" / "weekly_meeting_template.docx"

PPT_HEADING = re.compile(r"^(.+?)[（(]\s*汇报人\s*[：:]\s*(.+?)[）)]\s*$")
DATE_PATTERN = re.compile(r"(20\d{2})[年/\.\-](\d{1,2})[月/\.\-](\d{1,2})日?")
PROJECT_HEADING = re.compile(
    r"^[一二三四五六七八九十百零〇0-9]+[、.．]\s*(.+?)[（(]\s*汇报人\s*[：:]\s*(.+?)[）)]\s*$"
)

SECTION_LABELS = {
    "current": ("本周工作完成情况", "本周完成情况", "本周进度"),
    "next": ("下周研发计划", "下周工作计划", "下周计划"),
    "issues": ("本周问题", "问题、风险", "问题"),
}
SECTION_DISPLAY = {"current": "本周工作完成情况：", "next": "下周研发计划：", "issues": "本周问题："}
IGNORED_SLIDE_TITLES = (
    "在建项目整体情况",
    "项目整体情况",
    "本周jira情况",
    "jira情况",
)

TITLE_ALIASES = {
    "南京银行": "南京银行混沌测试",
    "南京银行项目周报": "南京银行混沌测试",
    "南京混沌测试": "南京银行混沌测试",
    "应急平台": "应急指挥调度平台",
    "应急演练指挥调度平台": "应急指挥调度平台",
    "中信证劵应急指挥调度平台二期": "中信证券应急指挥调度平台二期",
    "国信指标平台": "国信证券指标管理平台",
    "五矿证券": "五矿证券日志管理项目",
    "五矿日志项目": "五矿证券日志管理项目",
    "墨巡miciusops智能运维平台": "墨巡MiciusOps智能运维平台",
    "市场工作成果与计划": "本周市场工作成果与计划",
}


def _notify(callback: ProgressCallback, stage: str, percent: int, detail: str) -> None:
    if callback:
        callback({"stage": stage, "percent": max(0, min(100, percent)), "detail": detail})


def _text(value: Any) -> str:
    return str(value or "").replace("\u3000", " ").strip()


def _normalize_title(value: str) -> str:
    text = PPT_HEADING.sub(r"\1", _text(value))
    text = re.sub(r"[\s:：,，、()（）\[\]【】\-—_/]+", "", text).lower()
    text = re.sub(r"(周报|项目报告|项目)$", "", text)
    return text


NORMALIZED_ALIASES = {
    _normalize_title(alias): _normalize_title(canonical)
    for alias, canonical in TITLE_ALIASES.items()
}


def _canonical_title(value: str) -> str:
    normalized = _normalize_title(value)
    return NORMALIZED_ALIASES.get(normalized, normalized)


def _title_score(left: str, right: str) -> tuple[float, str]:
    left_raw = _normalize_title(left)
    right_raw = _normalize_title(right)
    if not left_raw or not right_raw:
        return 0.0, "none"
    if left_raw == right_raw:
        return 1.0, "exact"
    left_key = _canonical_title(left)
    right_key = _canonical_title(right)
    if left_key == right_key:
        return 0.96, "alias"
    shorter, longer = sorted((left_key, right_key), key=len)
    if len(shorter) >= 4 and shorter in longer and len(shorter) / len(longer) >= 0.55:
        return 0.84, "similar"
    return SequenceMatcher(None, left_key, right_key).ratio(), "similar"


def _clean_line(value: str) -> str:
    return re.sub(r"^[\s·•▪◆◇■□]+", "", _text(value)).strip()


def _clean_content(value: str) -> str:
    lines = []
    for raw in re.split(r"[\r\n\v]+", _text(value)):
        line = _clean_line(raw)
        if line and not re.fullmatch(r"[XxＸ]+(?:【.*?】)?", line):
            lines.append(line)
    return "\n".join(dict.fromkeys(lines))


def _date_value(value: str) -> str:
    match = DATE_PATTERN.search(_text(value))
    if not match:
        return _text(value)
    return f"{int(match.group(1)):04d}-{int(match.group(2)):02d}-{int(match.group(3)):02d}"


def current_week_saturday() -> date:
    """按周一至周日计算本周周六，周日仍归入刚结束的这一周。"""
    today = date.today()
    monday = today - timedelta(days=today.weekday())
    return monday + timedelta(days=5)


def _unique_cells(row) -> list[Any]:
    cells, seen = [], set()
    for cell in row.cells:
        marker = id(cell._tc)
        if marker not in seen:
            cells.append(cell)
            seen.add(marker)
    return cells


def _cell_lines(cell) -> list[str]:
    lines = []
    for paragraph in cell.paragraphs:
        lines.extend(part for part in re.split(r"[\r\n\v]+", paragraph.text) if part.strip())
    return lines


def _meeting_cell(document: Document):
    for table in document.tables:
        for row in table.rows:
            cells = _unique_cells(row)
            if cells and re.sub(r"\s+", "", cells[0].text) == "会议主要内容":
                return cells[-1]
    raise ValueError("部门周例会模板中没有找到“会议主要内容”区域")


def _meeting_date_cell(document: Document):
    for table in document.tables:
        for row in table.rows:
            cells = _unique_cells(row)
            if len(cells) >= 2 and re.sub(r"\s+", "", cells[0].text) == "会议时间":
                return cells[1]
    raise ValueError("部门周例会模板中没有找到“会议时间”字段")


@lru_cache(maxsize=1)
def _template_dimensions() -> tuple[int, int]:
    """缓存模板页面尺寸，避免每个源文件都重复打开模板。"""
    presentation = Presentation(PPT_TEMPLATE)
    return int(presentation.slide_width), int(presentation.slide_height)


def _project_slide_numbers(presentation: Presentation) -> list[int]:
    return [
        slide_number
        for slide_number, slide in enumerate(presentation.slides, start=1)
        if any(
            entry["text"] and PPT_HEADING.match(entry["text"])
            for entry in _flatten_shapes(slide.shapes)
        )
    ]


def _template_projects() -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    """读取内置 PPT 和 Word 模板中的项目顺序，不把项目名单硬编码在处理器内。"""
    ppt = Presentation(PPT_TEMPLATE)
    deck_projects = []
    seen = set()
    for slide_number, slide in enumerate(ppt.slides, start=1):
        title = next((
            entry["text"] for entry in _flatten_shapes(slide.shapes)
            if entry["text"] and PPT_HEADING.match(entry["text"])
        ), "")
        match = PPT_HEADING.match(title)
        if not match:
            continue
        key = _canonical_title(match.group(1))
        if key in seen:
            continue
        seen.add(key)
        deck_projects.append({"id": f"deck_{len(deck_projects) + 1}", "key": key, "title": match.group(1), "reporter": match.group(2), "template_slide": slide_number})

    document = Document(DOCX_TEMPLATE)
    meeting_projects = []
    for paragraph in _meeting_cell(document).paragraphs:
        match = PROJECT_HEADING.match(_text(paragraph.text))
        if not match:
            continue
        key = _canonical_title(match.group(1))
        meeting_projects.append({"id": f"meeting_{len(meeting_projects) + 1}", "key": key, "title": match.group(1), "reporter": match.group(2)})
    return deck_projects, meeting_projects


def _flatten_shapes(shapes, parent_left: int = 0, parent_top: int = 0, prefix: str = "") -> list[dict[str, Any]]:
    entries = []
    for index, shape in enumerate(shapes):
        path = f"{prefix}{index + 1}"
        left = parent_left + int(shape.left)
        top = parent_top + int(shape.top)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            entries.extend(_flatten_shapes(shape.shapes, left, top, f"{path}."))
            continue
        text = _text(getattr(shape, "text", ""))
        entries.append({
            "path": path,
            "shape": shape,
            "text": text,
            "left": left,
            "top": top,
            "width": int(shape.width),
            "height": int(shape.height),
        })
    return entries


def _section_label(text: str, key: str | None = None) -> str:
    labels = SECTION_LABELS[key] if key else tuple(label for values in SECTION_LABELS.values() for label in values)
    value = _text(text)
    compact = re.sub(r"\s+", "", value)
    first_line = re.sub(r"\s+", "", re.split(r"[\r\n\v]+", value, maxsplit=1)[0])
    for label in sorted(labels, key=len, reverse=True):
        if compact == label or compact.startswith((f"{label}：", f"{label}:")):
            return label
        if first_line in {label, f"{label}：", f"{label}:"}:
            return label
    return ""


def _is_section_label(text: str, key: str | None = None) -> bool:
    return bool(_section_label(text, key))


def _content_after_section_label(text: str, key: str) -> str:
    label = _section_label(text, key)
    if not label:
        return ""
    spaced_label = r"\s*".join(map(re.escape, label))
    remainder = re.sub(rf"^\s*{spaced_label}\s*[：:]?\s*", "", _text(text), count=1)
    return _clean_content(remainder)


def _section_from_entries(entries: list[dict[str, Any]], key: str) -> str:
    label_entries = [item for item in entries if _is_section_label(item["text"], key)]
    if not label_entries:
        return ""
    label = min(label_entries, key=lambda item: (item["top"], item["left"]))
    next_tops = [
        item["top"] for item in entries
        if item is not label and item["top"] > label["top"] + 1000 and _is_section_label(item["text"])
    ]
    bottom = min(next_tops) if next_tops else 10**10
    values = []
    inline_content = _content_after_section_label(label["text"], key)
    if inline_content:
        values.append((label["top"], label["left"], inline_content))
    for item in entries:
        if item is label or item["top"] < label["top"] - 100000 or item["top"] >= bottom:
            continue
        if item["left"] + 100000 < label["left"]:
            continue
        if key in {"current", "next"} and label["left"] < 2000000 and item["left"] > 7200000:
            continue
        if key == "issues" and label["left"] > 6000000 and item["left"] < 6000000:
            continue
        if _is_section_label(item["text"]) or PPT_HEADING.match(item["text"]):
            continue
        content = _clean_content(item["text"])
        if content:
            values.append((item["top"], item["left"], content))
    return "\n".join(dict.fromkeys(value for _, _, value in sorted(values)))


def _slide_title(entries: list[dict[str, Any]]) -> tuple[str, str, str]:
    for item in entries:
        match = PPT_HEADING.match(item["text"])
        if match:
            return match.group(1).strip(), match.group(2).strip(), "explicit"
    candidates = [item for item in entries if item["top"] < 1200000 and item["text"] and not _is_section_label(item["text"])]
    if not candidates:
        return "", "", "missing"
    return min(candidates, key=lambda item: (item["top"], item["left"]))["text"], "", "inferred"


def _is_ignored_slide(texts: list[str], slide_number: int, slide_count: int) -> bool:
    """识别不参与项目核算的封面、结束页和固定汇总页。"""
    if slide_number in {1, slide_count}:
        return True
    compact_texts = [re.sub(r"\s+", "", _text(text)).lower() for text in texts]
    return any(
        title in text
        for text in compact_texts
        for title in IGNORED_SLIDE_TITLES
    )


def _audit_shape_bounds(entries: list[dict[str, Any]], width: int, height: int, file: str, slide: int, project: str) -> list[dict[str, Any]]:
    issues = []
    for entry in entries:
        if entry["left"] < 0 or entry["top"] < 0 or entry["left"] + entry["width"] > width or entry["top"] + entry["height"] > height:
            issues.append({
                "severity": "warning", "code": "shape_overflow", "label": "对象超出画布",
                "file": file, "slide": slide, "location": f"对象 {entry['path']}", "project": project,
                "detail": "对象边界超出 PPT 页面，可能导致导出后裁切。",
                "suggestion": "回到原 PPT 调整对象大小或位置。",
            })
    return issues


def parse_presentation_source(path: str | Path, display_name: str, expected: list[dict[str, Any]]) -> dict[str, Any]:
    presentation = Presentation(path)
    width, height = int(presentation.slide_width), int(presentation.slide_height)
    template_width, template_height = _template_dimensions()
    issues = []
    if (width, height) != (template_width, template_height):
        issues.append({"severity": "error", "code": "slide_size", "label": "页面尺寸不一致", "file": display_name, "slide": 0, "location": "整份 PPT", "project": "", "detail": f"当前尺寸 {width}x{height}，模板尺寸 {template_width}x{template_height}。", "suggestion": "使用部门项目周报模板的宽高。"})
    slides = []
    slide_count = len(presentation.slides)
    for slide_number, slide in enumerate(presentation.slides, start=1):
        entries = _flatten_shapes(slide.shapes)
        texts = [entry["text"] for entry in entries if entry["text"]]
        has_sections = any(_is_section_label(text) for text in texts)
        if not has_sections and _is_ignored_slide(texts, slide_number, slide_count):
            continue
        if not texts:
            issues.append({"severity": "warning", "code": "blank_slide", "label": "空白页", "file": display_name, "slide": slide_number, "location": "整页", "project": "", "detail": "页面没有可识别文字内容。", "suggestion": "确认是否是误上传的空白页。"})
            continue
        if not has_sections:
            issues.append({"severity": "info", "code": "non_project_slide", "label": "非项目页", "file": display_name, "slide": slide_number, "location": "整页", "project": "", "detail": "封面、结束页或说明页未纳入项目合并。", "suggestion": "如需纳入，请在模板中增加对应项目页。"})
            continue
        title, reporter, title_mode = _slide_title(entries)
        if not title:
            issues.append({"severity": "error", "code": "title_missing", "label": "项目标题缺失", "file": display_name, "slide": slide_number, "location": "标题区域", "project": "", "detail": "项目页存在内容字段，但没有可识别的项目标题。", "suggestion": "补充项目名称并放在页面顶部。"})
            continue
        scores = sorted((_title_score(title, item["title"]) + (item,) for item in expected), key=lambda value: value[0], reverse=True)
        best_score, best_mode, best_project = scores[0] if scores else (0, "none", {"key": "", "title": ""})
        if best_score < 0.58:
            best_project = {"key": "", "title": ""}
        project_title = best_project.get("title", "")
        if (best_mode in {"alias", "similar"} or title_mode == "inferred") and project_title:
            issues.append({"severity": "warning", "code": "title_alias", "label": "标题需确认", "file": display_name, "slide": slide_number, "location": f"对象 {next((e['path'] for e in entries if e['text'] == title), '标题区域')}", "project": project_title, "detail": f"源标题“{title}”按别名或相似规则对应“{project_title}”。", "suggestion": "确认该页确实属于对应项目。"})
        if not reporter and project_title:
            issues.append({"severity": "warning", "code": "reporter_missing", "label": "汇报人缺失", "file": display_name, "slide": slide_number, "location": "标题区域", "project": project_title, "detail": "标题没有包含“汇报人”信息。", "suggestion": "按模板补充汇报人。"})
        elif reporter and project_title and _normalize_title(reporter) != _normalize_title(best_project.get("reporter", "")):
            issues.append({"severity": "warning", "code": "reporter_mismatch", "label": "汇报人不一致", "file": display_name, "slide": slide_number, "location": "标题区域", "project": project_title, "detail": f"源文件为“{reporter}”，模板为“{best_project.get('reporter', '')}”。", "suggestion": "以模板中的汇报人为准，并在审核结果确认。"})
        placeholders = [entry for entry in entries if re.search(r"X{2,}|Ｘ{2,}", entry["text"])]
        for entry in placeholders:
            issues.append({"severity": "error", "code": "placeholder", "label": "模板占位符未替换", "file": display_name, "slide": slide_number, "location": f"对象 {entry['path']}", "project": project_title, "detail": f"发现未替换内容：{entry['text'][:80]}", "suggestion": "补充真实周报内容后再合并。"})
        section_values = {
            section_key: _section_from_entries(entries, section_key)
            for section_key in SECTION_DISPLAY
        }
        for section_key, section_label in SECTION_DISPLAY.items():
            section_exists = any(_is_section_label(entry["text"], section_key) for entry in entries)
            if not section_exists or (section_key != "issues" and not section_values[section_key]):
                issues.append({
                    "severity": "warning",
                    "code": "section_missing",
                    "label": "内容字段缺失",
                    "file": display_name,
                    "slide": slide_number,
                    "location": f"{section_label}区域",
                    "project": project_title,
                    "detail": (
                        f"未识别到“{section_label.rstrip('：')}”字段。"
                        if not section_exists
                        else f"“{section_label.rstrip('：')}”字段没有有效内容。"
                    ),
                    "suggestion": "按项目周报模板补充该字段；没有问题时填写“无”。",
                })
        issues.extend(_audit_shape_bounds(entries, width, height, display_name, slide_number, project_title))
        slides.append({
            "id": f"{display_name}#{slide_number}", "file": display_name, "slide": slide_number,
            "title": title, "reporter": reporter, "title_mode": title_mode,
            "project_key": best_project.get("key", "") if project_title else "",
            "project_title": project_title, "score": round(best_score, 3),
            "current": section_values["current"], "next": section_values["next"], "issues": section_values["issues"],
        })
    return {"file": display_name, "slide_count": len(presentation.slides), "slides": slides, "issues": issues, "size": {"width": width, "height": height}}


def _is_generated_artifact(display_name: str) -> bool:
    stem = Path(display_name.replace("\\", "/")).stem.strip()
    return bool(re.match(r"^(项目周报|部门周例会)\s*\d", stem))


def _project_result(project: dict[str, Any], slides: list[dict[str, Any]], issues: list[dict[str, Any]], source_kind: str) -> dict[str, Any]:
    current = "\n".join(dict.fromkeys(item["current"] for item in slides if item["current"]))
    next_plan = "\n".join(dict.fromkeys(item["next"] for item in slides if item["next"]))
    problems = "\n".join(dict.fromkeys(item["issues"] for item in slides if item["issues"])) or "无"
    project_issues = [item for item in issues if item.get("project") == project["title"] or item.get("project") == project.get("key")]
    has_error = any(item["severity"] == "error" for item in project_issues)
    has_warning = any(item["severity"] == "warning" for item in project_issues)
    status = "需处理" if has_error or not slides else "待确认" if has_warning else "通过"
    if not slides:
        project_issues.append({"severity": "error", "code": "project_missing", "label": "项目页缺失", "file": "", "slide": 0, "location": "项目清单", "project": project["title"], "detail": "压缩包中没有找到该项目的有效内容页。", "suggestion": "补充对应项目周报 PPTX。"})
        status = "需处理"
    return {
        "id": project["id"], "key": project["key"], "title": project["title"], "reporter": project["reporter"],
        "template_slide": project.get("template_slide"),
        "status": status, "source_kind": source_kind, "slides": slides, "source_files": sorted({item["file"] for item in slides}),
        "current": current, "next": next_plan, "issues": problems, "checks": project_issues,
    }


def process_weekly_report(
    presentation_sources: list[tuple[Path, str]],
    source_manifest: list[dict[str, Any]],
    progress_callback: ProgressCallback = None,
) -> dict[str, Any]:
    """以 ZIP 内 PPTX 为内容源，对照内置 PPT/Word 模板完成审核和合并计划。"""
    deck_projects, meeting_projects = _template_projects()
    expected = []
    for item in deck_projects + meeting_projects:
        if not any(existing["key"] == item["key"] for existing in expected):
            expected.append(item)
    _notify(progress_callback, "读取模板", 10, f"已载入 {len(deck_projects)} 个项目 PPT 页模板和 {len(meeting_projects)} 个会议项目")

    parsed_files = []
    all_issues = []
    for index, (path, display_name) in enumerate(presentation_sources, start=1):
        if _is_generated_artifact(display_name):
            artifact_issue = {"severity": "info", "code": "generated_artifact", "label": "历史成品已忽略", "file": display_name, "slide": 0, "location": "文件", "project": "", "detail": "该文件看起来是历史汇总 PPT，不作为本次项目源。", "suggestion": "保留在目录中用于追溯，但不参与合并。"}
            parsed_files.append({"file": display_name, "slide_count": 0, "slides": [], "issues": [artifact_issue]})
            all_issues.append(artifact_issue)
            continue
        parsed = parse_presentation_source(path, display_name, expected)
        parsed_files.append(parsed)
        all_issues.extend(parsed["issues"])
        _notify(progress_callback, "审核项目 PPT", 12 + round(index / max(len(presentation_sources), 1) * 52), f"已审核 {index} / {len(presentation_sources)} 个文件")

    all_slides = [slide for parsed in parsed_files for slide in parsed["slides"]]
    by_key: dict[str, list[dict[str, Any]]] = {item["key"]: [] for item in expected}
    for slide in all_slides:
        if slide["project_key"]:
            by_key.setdefault(slide["project_key"], []).append(slide)
        else:
            all_issues.append({"severity": "warning", "code": "unmatched_slide", "label": "项目页未对应", "file": slide["file"], "slide": slide["slide"], "location": "标题区域", "project": "", "detail": f"源标题“{slide['title']}”未匹配到模板项目。", "suggestion": "确认是否应新增模板项目，或修正源 PPT 标题。"})

    deck_results = [_project_result(item, by_key.get(item["key"], []), all_issues, "项目周报") for item in deck_projects]
    meeting_results = [_project_result(item, by_key.get(item["key"], []), all_issues, "周例会") for item in meeting_projects]
    existing_missing = {(item.get("project"), item.get("code")) for item in all_issues}
    for project in deck_results + meeting_results:
        for issue in project["checks"]:
            marker = (issue.get("project"), issue.get("code"))
            if issue.get("code") == "project_missing" and marker not in existing_missing:
                all_issues.append(issue)
                existing_missing.add(marker)
    for project in deck_results + meeting_results:
        if len(project["source_files"]) > 1:
            all_issues.append({"severity": "warning", "code": "multiple_sources", "label": "多个源文件合并", "file": "、".join(project["source_files"]), "slide": 0, "location": project["title"], "project": project["title"], "detail": "同一项目来自多个 PPT 文件，系统按源文件和页码顺序合并。", "suggestion": "确认这些文件是否是同一项目的不同内容页。"})

    if not any(slide["project_key"] == _canonical_title("本周市场工作成果与计划") for slide in all_slides):
        all_issues.append({"severity": "warning", "code": "market_missing", "label": "市场周报缺失", "file": "", "slide": 0, "location": "部门周例会项目清单", "project": "本周市场工作成果与计划", "detail": "周例会模板包含市场工作项目，但源 PPT 中没有对应内容页。", "suggestion": "补充张珂珂的市场工作周报。"})

    for parsed in parsed_files:
        if parsed["file"] and not parsed["slides"] and not _is_generated_artifact(parsed["file"]):
            all_issues.append({"severity": "warning", "code": "no_project_slides", "label": "未识别项目页", "file": parsed["file"], "slide": 0, "location": "整份 PPT", "project": "", "detail": "文件中没有包含项目字段的有效页面。", "suggestion": "检查是否为封面、结束页或不符合模板的版式。"})

    week_end = current_week_saturday()
    error_count = sum(item["severity"] == "error" for item in all_issues)
    warning_count = sum(item["severity"] == "warning" for item in all_issues)
    project_count = len(deck_results)
    result = {
        "ok": True,
        "week_end": week_end.isoformat(),
        "output_stem": f"项目周报{week_end:%m%d}",
        "stats": {
            "project_count": project_count,
            "matched_projects": sum(bool(item["slides"]) for item in deck_results),
            "passed": sum(item["status"] == "通过" for item in deck_results),
            "pending": sum(item["status"] == "待确认" for item in deck_results),
            "issues": sum(item["status"] == "需处理" for item in deck_results),
            "error_count": error_count,
            "warning_count": warning_count,
            "source_files": len(presentation_sources),
            "source_slides": len(all_slides),
        },
        "projects": deck_results,
        "meeting_projects": meeting_results,
        "issues": all_issues,
        "files": parsed_files,
        "sources": source_manifest,
        "assembly": [slide for project in deck_results for slide in project["slides"]],
    }
    _notify(progress_callback, "完成审核", 88, f"项目页审核完成，发现 {error_count} 个错误、{warning_count} 个待确认项")
    _notify(progress_callback, "生成结果", 100, "可下载项目周报 PPT、周例会 DOCX 和审核报告")
    return result


def _remove_element(element) -> None:
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def _remove_slide(presentation: Presentation, slide) -> None:
    slide_index = next(index for index, candidate in enumerate(presentation.slides) if candidate.part is slide.part)
    slide_id = presentation.slides._sldIdLst[slide_index]
    presentation.part.drop_rel(slide_id.rId)
    presentation.slides._sldIdLst.remove(slide_id)


def _clear_paragraph_runs(paragraph, value: str) -> None:
    runs = paragraph.runs
    if runs:
        runs[0].text = value
        for run in runs[1:]:
            run.text = ""
    else:
        paragraph.text = value


def _copy_relationships(source_slide, target_slide, element) -> None:
    """复制图片/外链关系，避免跨 PPTX 的 rId 悬空。"""
    for node in element.iter():
        for attribute, source_rid in list(node.attrib.items()):
            if not attribute.startswith("{" + qn("r:id").split("}")[0].strip("{") + "}"):
                continue
            try:
                source_rel = source_slide.part.rels[source_rid]
            except KeyError:
                continue
            if source_rel.reltype == RT.IMAGE:
                _, target_rid = target_slide.part.get_or_add_image_part(BytesIO(source_rel.target_part.blob))
            elif source_rel.is_external:
                target_rid = target_slide.part.relate_to(source_rel.target_ref, source_rel.reltype, is_external=True)
            elif source_rel.reltype.endswith("/tags"):
                # 标签是 PowerPoint 的可选元数据，跨文件复制没有视觉影响。
                parent = node.getparent()
                if parent is not None:
                    parent.remove(node)
                continue
            else:
                raise ValueError(f"无法复制 {source_rel.reltype} 关系，请先在源 PPT 中展开对象")
            node.set(attribute, target_rid)


def _matching_layout(presentation: Presentation, source_slide):
    source_name = getattr(source_slide.slide_layout, "name", "")
    fallback = presentation.slide_layouts[-1]
    return next((layout for layout in presentation.slide_layouts if layout.name == source_name), fallback)


def _clone_source_slide(presentation: Presentation, source_slide):
    target_slide = presentation.slides.add_slide(_matching_layout(presentation, source_slide))
    for shape in list(target_slide.shapes):
        _remove_element(shape._element)
    for shape in source_slide.shapes:
        element = deepcopy(shape._element)
        _copy_relationships(source_slide, target_slide, element)
        target_slide.shapes._spTree.insert_element_before(element, "p:extLst")
    return target_slide


def _text_role(text: str) -> str:
    if PPT_HEADING.match(_text(text)):
        return "title"
    for key in SECTION_LABELS:
        if _is_section_label(text, key):
            return key
    return "body"


def _copy_xml_contents(target, source) -> None:
    target.attrib.clear()
    target.attrib.update(source.attrib)
    for child in list(target):
        target.remove(child)
    for child in source:
        target.append(deepcopy(child))


def _copy_text_style(source_shape, template_shape) -> None:
    """复制模板的段落与字体样式，保留源文本框几何和自动适配设置。"""
    source_frame = source_shape.text_frame
    template_frame = template_shape.text_frame
    source_paragraphs = source_frame.paragraphs
    template_paragraphs = template_frame.paragraphs
    if not template_paragraphs:
        return
    for paragraph_index, source_paragraph in enumerate(source_paragraphs):
        template_paragraph = template_paragraphs[min(paragraph_index, len(template_paragraphs) - 1)]
        template_ppr = template_paragraph._p.find(qn("a:pPr"))
        if template_ppr is not None:
            source_ppr = source_paragraph._p.get_or_add_pPr()
            _copy_xml_contents(source_ppr, template_ppr)
        else:
            source_ppr = source_paragraph._p.find(qn("a:pPr"))
            if source_ppr is not None:
                source_paragraph._p.remove(source_ppr)
        template_run = next(iter(template_paragraph.runs), None)
        if template_run is None:
            continue
        template_rpr = template_run._r.find(qn("a:rPr"))
        for source_run in source_paragraph.runs:
            if template_rpr is not None:
                source_rpr = source_run._r.get_or_add_rPr()
                _copy_xml_contents(source_rpr, template_rpr)
            else:
                source_rpr = source_run._r.find(qn("a:rPr"))
                if source_rpr is not None:
                    source_run._r.remove(source_rpr)


def _apply_template_text_format(source_slide, template_slide) -> None:
    source_entries = [entry for entry in _flatten_shapes(source_slide.shapes) if getattr(entry["shape"], "has_text_frame", False) and entry["text"]]
    template_entries = [entry for entry in _flatten_shapes(template_slide.shapes) if getattr(entry["shape"], "has_text_frame", False)]
    if not source_entries or not template_entries:
        return
    for source_entry in source_entries:
        template_entry = _matching_template_entry(source_entry, template_entries)
        if template_entry is None:
            continue
        _copy_text_style(source_entry["shape"], template_entry["shape"])


def _shape_font_size(shape, default: float = 12.0) -> float:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                return max(6.0, float(run.font.size.pt))
    return default


def _shape_capacity(shape) -> tuple[int, int]:
    """估算当前字号下文本框可容纳的字符数和行数。"""
    font_size = _shape_font_size(shape)
    emu_per_point = 12700
    chars_per_line = max(8, int(int(shape.width) / (font_size * emu_per_point * 0.95)))
    max_lines = max(1, int(int(shape.height) / (font_size * emu_per_point * 1.3)))
    return chars_per_line, max_lines


def _fit_body_font(shape, minimum: int = 9) -> None:
    text = _text(shape.text)
    if not text:
        return
    current_size = _shape_font_size(shape)
    start = max(minimum, int(round(current_size)))
    for size in range(start, minimum - 1, -1):
        chars_per_line, max_lines = _shape_capacity_for_font(shape, size)
        if len(text) <= chars_per_line * max_lines:
            if size < current_size:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(size)
            return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(minimum)


def _shape_capacity_for_font(shape, font_size: float) -> tuple[int, int]:
    emu_per_point = 12700
    chars_per_line = max(8, int(int(shape.width) / (font_size * emu_per_point * 0.95)))
    max_lines = max(1, int(int(shape.height) / (font_size * emu_per_point * 1.3)))
    return chars_per_line, max_lines


def _text_chunks(shape) -> list[str]:
    text = _text(shape.text)
    if not text:
        return [""]
    chars_per_line, max_lines = _shape_capacity(shape)
    lines = []
    for raw_line in text.splitlines() or [text]:
        line = raw_line or " "
        while len(line) > chars_per_line:
            lines.append(line[:chars_per_line])
            line = line[chars_per_line:]
        lines.append(line)
    chunk_size = max(1, max_lines)
    return ["\n".join(lines[index:index + chunk_size]) for index in range(0, len(lines), chunk_size)]


def _overflow_chunks(slide) -> dict[str, list[str]]:
    chunks = {}
    for entry in _flatten_shapes(slide.shapes):
        if not entry["text"] or _text_role(entry["text"]) != "body":
            continue
        _fit_body_font(entry["shape"])
        values = _text_chunks(entry["shape"])
        if len(values) > 1:
            chunks[entry["path"]] = values
    return chunks


def _set_page_chunks(slide, chunks: dict[str, list[str]], page_index: int) -> None:
    for entry in _flatten_shapes(slide.shapes):
        if not entry["text"] or _text_role(entry["text"]) != "body":
            continue
        values = chunks.get(entry["path"])
        if values is not None:
            entry["shape"].text = values[page_index] if page_index < len(values) else ""
        elif page_index:
            entry["shape"].text = ""


def _normalized_content(value: str) -> str:
    return re.sub(r"\s+", "", _text(value))


def _xml_signature(element, ignored_attributes: set[str] | None = None):
    if element is None:
        return None
    ignored = ignored_attributes or set()
    attributes = tuple(sorted(
        (name, value) for name, value in element.attrib.items()
        if name.rsplit("}", 1)[-1] not in ignored
    ))
    return element.tag, attributes, tuple(_xml_signature(child, ignored) for child in element)


def _text_style_matches(shape, template_shape) -> bool:
    template_paragraphs = template_shape.text_frame.paragraphs
    if not template_paragraphs:
        return True
    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
        template_paragraph = template_paragraphs[min(paragraph_index, len(template_paragraphs) - 1)]
        if _xml_signature(paragraph._p.find(qn("a:pPr"))) != _xml_signature(template_paragraph._p.find(qn("a:pPr"))):
            return False
        template_run = next(iter(template_paragraph.runs), None)
        template_rpr = template_run._r.find(qn("a:rPr")) if template_run is not None else None
        template_signature = _xml_signature(template_rpr, {"sz"})
        for run in paragraph.runs:
            if _xml_signature(run._r.find(qn("a:rPr")), {"sz"}) != template_signature:
                return False
    return True


def _matching_template_entry(source_entry: dict[str, Any], template_entries: list[dict[str, Any]]):
    role = _text_role(source_entry["text"])
    candidates = [entry for entry in template_entries if _text_role(entry["text"]) == role]
    if not candidates:
        candidates = [entry for entry in template_entries if _text_role(entry["text"]) == "body"]
    if not candidates:
        return None
    return min(
        candidates,
        key=lambda entry: abs(entry["left"] - source_entry["left"]) + abs(entry["top"] - source_entry["top"]),
    )


def _generation_qa(
    presentation: Presentation,
    generated_groups: list[dict[str, Any]],
    max_rounds: int = 4,
    progress_callback: ProgressCallback = None,
) -> dict[str, Any]:
    """多轮核验内容、模板文字样式和版面稳定性，并修复可确定的问题。"""
    history = []
    stable_rounds = 0
    final_issues = []
    thresholds = {"content": 0.99, "style": 0.97, "layout": 0.97}
    for round_number in range(1, max_rounds + 1):
        _notify(progress_callback, "生成质量核验", 91 + round_number, f"正在执行第 {round_number} / {max_rounds} 轮内容、样式与版面核验")
        issues = []
        repairs = 0
        content_total = content_ok = 0
        style_total = style_ok = 0
        layout_total = layout_ok = 0
        for group in generated_groups:
            for slide in group["slides"]:
                content_total += 1
                title, reporter, _ = _slide_title(_flatten_shapes(slide.shapes))
                if _canonical_title(title) == group["project"]["key"] and _normalize_title(reporter) == _normalize_title(group["project"]["reporter"]):
                    content_ok += 1
                else:
                    _replace_title(slide, group["project"]["title"], group["project"]["reporter"])
                    repairs += 1
            source_entries = {
                entry["path"]: entry for entry in _flatten_shapes(group["source_slide"].shapes)
                if entry["text"] and _text_role(entry["text"]) == "body"
            }
            page_maps = [
                {entry["path"]: entry for entry in _flatten_shapes(slide.shapes)}
                for slide in group["slides"]
            ]
            for path, source_entry in source_entries.items():
                content_total += 1
                expected = _normalized_content(source_entry["text"])
                actual = _normalized_content("".join(
                    page[path]["text"] for page in page_maps
                    if path in page and page[path]["text"]
                ))
                if actual == expected:
                    content_ok += 1
                else:
                    issues.append({
                        "severity": "error", "code": "generated_content_mismatch", "label": "生成内容不一致",
                        "file": group["source"]["file"], "slide": group["source"]["slide"],
                        "location": f"对象 {path}", "project": group["project"]["title"],
                        "detail": "生成页中的正文与源 PPT 不一致，可能存在重复或缺失。",
                        "suggestion": "系统已重新写入分页内容；请在最终审核结果中确认。",
                    })
                    for page_index, slide in enumerate(group["slides"]):
                        _set_page_chunks(slide, group["chunks"], page_index)
                    repairs += 1

            template_slide = group.get("template_slide")
            template_entries = [
                entry for entry in _flatten_shapes(template_slide.shapes)
                if getattr(entry["shape"], "has_text_frame", False)
            ] if template_slide is not None else []
            for slide in group["slides"]:
                for entry in _flatten_shapes(slide.shapes):
                    shape = entry["shape"]
                    if entry["text"] and getattr(shape, "has_text_frame", False):
                        layout_total += 1
                        if len(_text_chunks(shape)) <= 1:
                            layout_ok += 1
                        else:
                            issues.append({
                                "severity": "warning", "code": "generated_text_overflow", "label": "生成文本仍可能溢出",
                                "file": group["source"]["file"], "slide": group["source"]["slide"],
                                "location": f"对象 {entry['path']}", "project": group["project"]["title"],
                                "detail": "自动缩放和分页后，文本量仍超过当前文本框估算容量。",
                                "suggestion": "建议人工检查该页文本框实际显示效果。",
                            })
                        template_entry = _matching_template_entry(entry, template_entries) if template_entries else None
                        if template_entry is not None:
                            style_total += 1
                            if _text_style_matches(shape, template_entry["shape"]):
                                style_ok += 1
                            else:
                                _copy_text_style(shape, template_entry["shape"])
                                if _text_role(entry["text"]) == "body":
                                    _fit_body_font(shape)
                                repairs += 1
                    layout_total += 1
                    if entry["left"] >= 0 and entry["top"] >= 0 and entry["left"] + entry["width"] <= presentation.slide_width and entry["top"] + entry["height"] <= presentation.slide_height:
                        layout_ok += 1

        scores = {
            "content": content_ok / content_total if content_total else 1.0,
            "style": style_ok / style_total if style_total else 1.0,
            "layout": layout_ok / layout_total if layout_total else 1.0,
        }
        passed = all(scores[key] >= value for key, value in thresholds.items())
        stable_rounds = stable_rounds + 1 if passed and repairs == 0 else 0
        history.append({
            "round": round_number,
            "content_score": round(scores["content"] * 100, 1),
            "style_score": round(scores["style"] * 100, 1),
            "layout_score": round(scores["layout"] * 100, 1),
            "repairs": repairs,
            "issues": len(issues),
            "stable": stable_rounds >= 2,
        })
        final_issues = issues
        if stable_rounds >= 2:
            break
    last = history[-1]
    overall_score = round(min(last["content_score"], last["style_score"], last["layout_score"]), 1)
    return {
        "status": "稳定" if stable_rounds >= 2 else "需人工复核",
        "stable": stable_rounds >= 2,
        "score": overall_score,
        "thresholds": {key: round(value * 100) for key, value in thresholds.items()},
        "rounds": history,
        "issues": final_issues,
    }


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _replace_title(slide, title: str, reporter: str) -> None:
    expected = f"{title}（汇报人：{reporter or '待补充'}）"
    candidates = [shape for shape in _iter_shapes(slide.shapes) if getattr(shape, "has_text_frame", False) and _text(shape.text)]
    title_shape = next((shape for shape in candidates if PPT_HEADING.match(_text(shape.text))), candidates[0] if candidates else None)
    if not title_shape:
        return
    paragraph = title_shape.text_frame.paragraphs[0]
    _clear_paragraph_runs(paragraph, expected)
    for extra in title_shape.text_frame.paragraphs[1:]:
        _clear_paragraph_runs(extra, "")


def build_weekly_presentation(
    result: dict[str, Any],
    source_lookup: dict[str, Path],
    target: str | Path,
    template: str | Path = PPT_TEMPLATE,
    progress_callback: ProgressCallback = None,
) -> None:
    """保留模板前三页和结束页，按源 PPT 页码顺序复制项目页。"""
    presentation = Presentation(template)
    week_end = datetime.strptime(result["week_end"], "%Y-%m-%d").date()
    cover_text = week_end.strftime("%Y年%m月%d日")
    for shape in presentation.slides[0].shapes:
        if getattr(shape, "has_text_frame", False) and DATE_PATTERN.search(_text(shape.text)):
            paragraph = shape.text_frame.paragraphs[0]
            _clear_paragraph_runs(paragraph, DATE_PATTERN.sub(cover_text, _text(shape.text)))

    placeholder_numbers = sorted(_project_slide_numbers(presentation), reverse=True)
    for slide_number in placeholder_numbers:
        _remove_slide(presentation, presentation.slides[slide_number - 1])
    outro_index = len(presentation.slides) - 1
    source_cache: dict[str, Presentation] = {}
    style_presentation = Presentation(template)
    style_slides = {
        project["key"]: style_presentation.slides[project["template_slide"] - 1]
        for project in result["projects"]
        if project.get("template_slide")
    }
    generated_groups = []

    def place_before_outro(slide) -> None:
        nonlocal outro_index
        slide_ids = presentation.slides._sldIdLst
        last = slide_ids[-1]
        slide_ids.remove(last)
        slide_ids.insert(outro_index, last)
        outro_index += 1

    for project in result["projects"]:
        for source in project["slides"]:
            display_name = source["file"]
            if display_name not in source_cache:
                source_cache[display_name] = Presentation(source_lookup[display_name])
            source_slide = source_cache[display_name].slides[source["slide"] - 1]
            cloned = _clone_source_slide(presentation, source_slide)
            template_slide = style_slides.get(project["key"])
            if template_slide is not None:
                _apply_template_text_format(cloned, template_slide)
            _replace_title(cloned, project["title"], project["reporter"])
            chunks = _overflow_chunks(cloned)
            page_count = max((len(values) for values in chunks.values()), default=1)
            generated_pages = []
            for page_index in range(page_count):
                if page_index:
                    continuation = _clone_source_slide(presentation, source_slide)
                    if template_slide is not None:
                        _apply_template_text_format(continuation, template_slide)
                    _replace_title(continuation, project["title"], project["reporter"])
                    cloned = continuation
                _set_page_chunks(cloned, chunks, page_index)
                if template_slide is not None:
                    _apply_template_text_format(cloned, template_slide)
                for entry in _flatten_shapes(cloned.shapes):
                    if entry["text"] and _text_role(entry["text"]) == "body":
                        _fit_body_font(entry["shape"])
                place_before_outro(cloned)
                generated_pages.append(cloned)
            generated_groups.append({
                "project": project,
                "source": source,
                "source_slide": source_slide,
                "template_slide": template_slide,
                "chunks": chunks,
                "slides": generated_pages,
            })
    qa = _generation_qa(presentation, generated_groups, progress_callback=progress_callback)
    result["qa"] = qa
    result["issues"].extend(qa["issues"])
    result["stats"]["qa_score"] = qa["score"]
    result["stats"]["qa_status"] = qa["status"]
    result["stats"]["error_count"] = sum(item["severity"] == "error" for item in result["issues"])
    result["stats"]["warning_count"] = sum(item["severity"] == "warning" for item in result["issues"])
    presentation.save(target)


def _set_paragraph_xml_text(paragraph_element, value: str) -> None:
    text_nodes = paragraph_element.xpath(".//w:t")
    if text_nodes:
        text_nodes[0].text = value
        for node in text_nodes[1:]:
            node.text = ""
    else:
        run = paragraph_element.find(qn("w:r"))
        if run is None:
            run = paragraph_element.makeelement(qn("w:r"), {})
            paragraph_element.append(run)
        text_node = run.find(qn("w:t"))
        if text_node is None:
            text_node = run.makeelement(qn("w:t"), {})
            run.append(text_node)
        text_node.text = value


def _paragraph_copy(template_paragraph, value: str):
    paragraph = deepcopy(template_paragraph._p)
    _set_paragraph_xml_text(paragraph, value)
    return paragraph


def build_weekly_meeting_document(result: dict[str, Any], target: str | Path, template: str | Path = DOCX_TEMPLATE) -> None:
    """在部门周例会模板的项目区域填充合并后 PPT 的三类内容。"""
    document = Document(template)
    date_cell = _meeting_date_cell(document)
    _clear_paragraph_runs(date_cell.paragraphs[0], datetime.strptime(result["week_end"], "%Y-%m-%d").strftime("%Y/%m/%d"))
    cell = _meeting_cell(document)
    source_paragraphs = list(cell.paragraphs)
    headings = [paragraph for paragraph in source_paragraphs if PROJECT_HEADING.match(_text(paragraph.text))]
    heading_by_key = {_canonical_title(PROJECT_HEADING.match(_text(paragraph.text)).group(1)): paragraph for paragraph in headings}
    label_templates = {}
    bullet_template = next((paragraph for paragraph in source_paragraphs if _clean_line(paragraph.text) in {"·", ""} and paragraph.text.strip() == "·"), headings[0])
    blank_template = next((paragraph for paragraph in source_paragraphs if not paragraph.text.strip()), bullet_template)
    for paragraph in source_paragraphs:
        compact = _text(paragraph.text).replace(" ", "")
        for key, labels in SECTION_LABELS.items():
            if any(compact.startswith(label) for label in labels):
                label_templates.setdefault(key, paragraph)

    tc = cell._tc
    for paragraph in list(cell.paragraphs):
        _remove_element(paragraph._p)
    for project in result["meeting_projects"]:
        heading = heading_by_key.get(project["key"]) or headings[0]
        tc.append(_paragraph_copy(heading, _text(heading.text)))
        for key in ("current", "next", "issues"):
            label_template = label_templates.get(key) or bullet_template
            tc.append(_paragraph_copy(label_template, SECTION_DISPLAY[key]))
            values = [line for line in project.get(key, "").splitlines() if line.strip()] or ["无"]
            for value in values:
                tc.append(_paragraph_copy(bullet_template, value))
            tc.append(deepcopy(blank_template._p))
    document.save(target)


def _excel_safe(value: Any) -> Any:
    if not isinstance(value, str):
        return value
    return "'" + value if value[:1] in {"=", "+", "-", "@"} else value


def export_weekly_report_xlsx(result: dict[str, Any], target: str | Path) -> None:
    workbook = Workbook()
    workbook.remove(workbook.active)
    header_fill = PatternFill("solid", fgColor="181816")
    header_font = Font(color="FFFFFF", bold=True)

    def add_sheet(title: str, headers: list[str], rows: list[list[Any]]) -> None:
        sheet = workbook.create_sheet(title)
        sheet.append(headers)
        for cell in sheet[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center")
        for row in rows:
            sheet.append([_excel_safe(value) for value in row])
        sheet.freeze_panes = "A2"
        sheet.auto_filter.ref = sheet.dimensions
        for column in sheet.columns:
            width = min(max(max(len(str(cell.value or "")) for cell in column) + 2, 10), 54)
            sheet.column_dimensions[column[0].column_letter].width = width
        for row in sheet.iter_rows(min_row=2):
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)

    add_sheet(
        "审核结果",
        ["严重程度", "文件", "页码", "位置", "项目", "检查项", "描述", "建议"],
        [[item["severity"], item["file"], item["slide"], item["location"], item["project"], item["label"], item["detail"], item["suggestion"]] for item in result["issues"]],
    )
    add_sheet(
        "项目整合",
        ["项目", "汇报人", "状态", "源文件", "源页码", "本周进度", "下周计划", "本周问题"],
        [[item["title"], item["reporter"], item["status"], "、".join(item["source_files"]), "、".join(str(slide["slide"]) for slide in item["slides"]), item["current"], item["next"], item["issues"]] for item in result["projects"]],
    )
    add_sheet(
        "周例会内容",
        ["项目", "汇报人", "本周完成", "下周计划", "本周问题"],
        [[item["title"], item["reporter"], item["current"], item["next"], item["issues"]] for item in result["meeting_projects"]],
    )
    add_sheet(
        "文件目录",
        ["路径", "类型", "大小", "状态"],
        [[item.get("path", ""), item.get("kind", ""), item.get("size", ""), item.get("status", "")] for item in result["sources"]],
    )
    qa = result.get("qa") or {}
    add_sheet(
        "生成质量核验",
        ["轮次", "内容准确度", "模板样式准确度", "版面稳定度", "自动修复数", "问题数", "连续稳定"],
        [[
            item.get("round", ""),
            item.get("content_score", ""),
            item.get("style_score", ""),
            item.get("layout_score", ""),
            item.get("repairs", ""),
            item.get("issues", ""),
            "是" if item.get("stable") else "否",
        ] for item in qa.get("rounds", [])],
    )
    workbook.save(target)
    workbook.close()
